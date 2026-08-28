import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation(json_file, output_pptx):
    with open(json_file, encoding='utf-8') as f:
        data = json.load(f)
        
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    # Colors
    DARK_BG = RGBColor(11, 15, 25)
    CARD_BG = RGBColor(30, 41, 59)
    CYAN = RGBColor(6, 182, 212)
    WHITE = RGBColor(255, 255, 255)
    MUTED = RGBColor(156, 163, 175)
    ACCENT = RGBColor(99, 102, 241)
    EMERALD = RGBColor(16, 185, 129)

    # 1. Title Slide
    slide = prs.slides.add_slide(blank_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    # Title Card Background
    shape = slide.shapes.add_shape(1, Inches(0.8), Inches(1.2), Inches(11.733), Inches(5.1)) # 1 is RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = ACCENT
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.8)
    tf.margin_top = Inches(0.8)

    p1 = tf.paragraphs[0]
    p1.text = "NCERT CLASS 11 CHEMISTRY • UNIT 4"
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = CYAN

    p2 = tf.add_paragraph()
    p2.text = "Chemical Bonding &\nMolecular Structure"
    p2.font.size = Pt(40)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.space_before = Pt(14)
    p2.space_after = Pt(20)

    p3 = tf.add_paragraph()
    p3.text = "Complete Visual Concept Library, VSEPR Geometries & Solved Exercises Suite"
    p3.font.size = Pt(18)
    p3.font.color.rgb = MUTED

    # 2. Module Slides
    for mod in data.get("modules", []):
        slide = prs.slides.add_slide(blank_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG

        # Slide Header Title
        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8))
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = mod["title"].upper()
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = CYAN

        # Subtitle / Summary Banner
        banner = slide.shapes.add_shape(1, Inches(0.8), Inches(1.3), Inches(11.733), Inches(1.1))
        banner.fill.solid()
        banner.fill.fore_color.rgb = CARD_BG
        banner.line.color.rgb = ACCENT

        btf = banner.text_frame
        btf.word_wrap = True
        btf.margin_left = Inches(0.4)
        btf.margin_top = Inches(0.2)
        bp = btf.paragraphs[0]
        bp.text = mod.get("summary") or mod.get("subtitle") or ""
        bp.font.size = Pt(14)
        bp.font.color.rgb = WHITE

        # Content List
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.733), Inches(4.2))
        ctf = content_box.text_frame
        ctf.word_wrap = True

        if mod.get("subsections"):
            for idx, sub in enumerate(mod["subsections"]):
                if idx < 3:
                    sp1 = ctf.add_paragraph() if idx > 0 else ctf.paragraphs[0]
                    sp1.text = f"• {sub['title']}"
                    sp1.font.size = Pt(16)
                    sp1.font.bold = True
                    sp1.font.color.rgb = ACCENT
                    sp1.space_before = Pt(10)

                    sp2 = ctf.add_paragraph()
                    snippet = sub['content'].replace('\n', ' ')
                    if len(snippet) > 220:
                        snippet = snippet[:220] + "..."
                    sp2.text = f"   {snippet}"
                    sp2.font.size = Pt(13)
                    sp2.font.color.rgb = WHITE

        # Footer
        footer = slide.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(11.733), Inches(0.4))
        ftf = footer.text_frame
        fp = ftf.paragraphs[0]
        fp.text = "ChemBondStudio • NCERT Class 11 Chemistry Chapter 4"
        fp.font.size = Pt(10)
        fp.font.color.rgb = MUTED

    prs.save(output_pptx)
    print(f"Successfully generated PowerPoint presentation: {output_pptx}")

if __name__ == "__main__":
    create_presentation(r"c:\college\bonding_data.json", r"c:\college\Chemical_Bonding_and_Molecular_Structure.pptx")
